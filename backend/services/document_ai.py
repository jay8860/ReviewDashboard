import json
import os
import re
import urllib.error
import urllib.parse
import urllib.request
from typing import Optional, Tuple
from datetime import date

try:
    from docx import Document
except Exception:
    Document = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json",
    ".docx", ".pptx", ".pdf", ".xlsx"
}

CONTEXT_TOKEN_LIMIT = int(os.getenv("DOC_ANALYSIS_CONTEXT_TOKEN_LIMIT", "20000"))
# Approximation: 1 token ~= 4 chars. This caps extracted text size to fit requested context.
MAX_EXTRACT_CHARS = int(os.getenv("DOC_ANALYSIS_MAX_CHARS", str(CONTEXT_TOKEN_LIMIT * 4)))
MAX_OUTPUT_TOKENS = int(os.getenv("DOC_ANALYSIS_MAX_OUTPUT_TOKENS", str(CONTEXT_TOKEN_LIMIT)))
MAX_CONTINUATIONS = int(os.getenv("DOC_ANALYSIS_MAX_CONTINUATIONS", "2"))
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.5-flash")


class GeminiModelNotFoundError(RuntimeError):
    pass


def _normalize_model_name(name: str) -> str:
    raw = (name or "").strip()
    if not raw:
        return ""
    match = re.search(r"models/([^/:]+)", raw, flags=re.IGNORECASE)
    if match:
        raw = match.group(1)
    raw = raw.replace(":generateContent", "")
    raw = raw.strip().strip("/")
    return raw


def _read_text_file(path: str) -> str:
    with open(path, "rb") as f:
        raw = f.read()
    return raw.decode("utf-8", errors="ignore")


def _read_docx(path: str) -> str:
    if Document is None:
        raise RuntimeError("python-docx is not installed. Install dependency to analyze .docx files.")
    doc = Document(path)
    chunks = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            chunks.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            if any(cells):
                chunks.append(" | ".join(cells))

    return "\n".join(chunks)


def _read_pptx(path: str) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Install dependency to analyze .pptx files.")
    presentation = Presentation(path)
    chunks = []
    for slide_no, slide in enumerate(presentation.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            chunks.append(f"Slide {slide_no}:")
            chunks.extend(slide_lines)
    return "\n".join(chunks)


def _read_pdf(path: str) -> str:
    if PdfReader is None:
        raise RuntimeError("pypdf is not installed. Install dependency to analyze .pdf files.")
    reader = PdfReader(path)
    chunks = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            chunks.append(f"Page {i}:")
            chunks.append(text)
    return "\n".join(chunks)


def _read_xlsx(path: str) -> str:
    if openpyxl is None:
        raise RuntimeError("openpyxl is not installed. Install dependency to analyze .xlsx files.")
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    chunks = []
    for ws in wb.worksheets[:6]:
        chunks.append(f"Sheet: {ws.title}")
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            values = ["" if v is None else str(v).strip() for v in row]
            if any(values):
                chunks.append("\t".join(values))
            row_count += 1
            if row_count >= 350:
                chunks.append("...")
                break
    return "\n".join(chunks)


def extract_text_from_document(file_path: str, extension: str) -> Tuple[str, bool]:
    ext = extension.lower()

    if ext in {".txt", ".md", ".csv", ".json"}:
        text = _read_text_file(file_path)
    elif ext == ".docx":
        text = _read_docx(file_path)
    elif ext == ".pptx":
        text = _read_pptx(file_path)
    elif ext == ".pdf":
        text = _read_pdf(file_path)
    elif ext == ".xlsx":
        text = _read_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")

    text = (text or "").strip()
    if not text:
        raise ValueError("No readable text found in the uploaded document.")

    is_truncated = len(text) > MAX_EXTRACT_CHARS
    return text[:MAX_EXTRACT_CHARS], is_truncated


def _build_default_prompt(document_name: str, text: str) -> str:
    return f"""
You are a governance and policy review analyst. Analyze the following document in detail.

Document: {document_name}

Return your answer in Markdown using exactly these sections:

1) Executive Summary
- 6-10 crisp bullet points.

2) Agenda-wise Findings
- Group findings into inferred agenda themes (Agenda 1, Agenda 2, ...).
- Under each agenda, provide point-wise observations and evidence.

3) Risks / Gaps / Bottlenecks
- Prioritized bullets with severity (High/Medium/Low).

4) Action Recommendations (Table)
- Create a markdown table with columns:
  Action Point | Suggested Owner | Timeline | Expected Outcome | Priority

5) Questions For Next Meeting
- 5-10 pointed questions.

6) Suggested Minutes (Draft)
- Draft concise meeting minutes with decisions and follow-ups.

Rules:
- Be specific and practical.
- If a fact is uncertain, mark it as "Assumption".
- Keep output actionable for an administrative review meeting.

Document text:
{text}
""".strip()


def _build_custom_prompt(document_name: str, text: str, custom_prompt: str) -> str:
    return f"""
You are a governance document analyst.

User instruction:
{custom_prompt}

Document: {document_name}

Return a practical, structured answer in Markdown with bullets/tables wherever useful.
Mention assumptions explicitly where evidence is missing.

Document text:
{text}
""".strip()


def _build_continuation_prompt(partial_output: str) -> str:
    tail = (partial_output or "")[-6000:]
    return f"""
The previous response was cut short due to output token limit.
Continue from exactly where it ended.

Rules:
- Do not repeat any text already produced.
- Keep the same structure and formatting.
- Return only the continuation content in Markdown.

Tail of already generated response:
{tail}
""".strip()


def _output_token_candidates() -> list[int]:
    requested = max(1024, int(MAX_OUTPUT_TOKENS))
    defaults = [requested, 20000, 16000, 12000, 8192, 4096, 2048]
    candidates = []
    for value in defaults:
        if value <= requested and value not in candidates:
            candidates.append(value)
    if requested not in candidates:
        candidates.insert(0, requested)
    return candidates


def _gemini_generate_once(api_key: str, model_name: str, prompt: str, output_tokens: int) -> Tuple[str, str]:
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": output_tokens,
        },
    }

    url = (
        f"https://generativelanguage.googleapis.com/v1beta/models/"
        f"{urllib.parse.quote(model_name)}:generateContent?key={urllib.parse.quote(api_key)}"
    )
    req = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json"},
        method="POST",
    )

    with urllib.request.urlopen(req, timeout=180) as resp:
        body = resp.read().decode("utf-8")
    parsed = json.loads(body)

    candidates = parsed.get("candidates") or []
    if not candidates:
        feedback = parsed.get("promptFeedback") or {}
        block_reason = feedback.get("blockReason")
        if block_reason:
            raise RuntimeError(f"Gemini blocked the response: {block_reason}")
        raise RuntimeError("Gemini returned no candidates.")

    candidate = candidates[0] if isinstance(candidates[0], dict) else {}
    finish_reason = str(candidate.get("finishReason") or "").upper()
    parts = candidate.get("content", {}).get("parts", [])
    text_chunks = [p.get("text", "") for p in parts if isinstance(p, dict)]
    result = "\n".join([c for c in text_chunks if c]).strip()
    if not result:
        raise RuntimeError(f"Gemini returned an empty response (finishReason={finish_reason or 'UNKNOWN'}).")
    return result, finish_reason


def _extract_json_from_text(raw_text: str):
    text = (raw_text or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?\s*", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\s*```$", "", text)
        text = text.strip()
    try:
        return json.loads(text)
    except Exception:
        pass

    candidates = []
    arr_start = text.find("[")
    obj_start = text.find("{")
    if arr_start != -1:
        candidates.append(arr_start)
    if obj_start != -1:
        candidates.append(obj_start)
    if not candidates:
        raise ValueError("No JSON found in model response.")

    start = min(candidates)
    for open_char, close_char in [("{", "}"), ("[", "]")]:
        if text[start] != open_char:
            continue
        depth = 0
        in_str = False
        esc = False
        for idx in range(start, len(text)):
            ch = text[idx]
            if in_str:
                if esc:
                    esc = False
                elif ch == "\\":
                    esc = True
                elif ch == '"':
                    in_str = False
                continue
            if ch == '"':
                in_str = True
                continue
            if ch == open_char:
                depth += 1
            elif ch == close_char:
                depth -= 1
                if depth == 0:
                    snippet = text[start:idx + 1]
                    return json.loads(snippet)
    raise ValueError("Could not parse JSON payload from model response.")


def _normalize_priority(value: Optional[str]) -> str:
    val = str(value or "").strip().lower()
    if val in {"critical", "crit", "p0"}:
        return "Critical"
    if val in {"high", "p1"}:
        return "High"
    if val in {"low", "p3"}:
        return "Low"
    return "Normal"


def generate_task_suggestions_with_gemini(
    source_name: str,
    source_text: str,
    department_name: Optional[str] = None,
    focus_prompt: Optional[str] = None,
) -> list[dict]:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is not set. Please add it in environment variables.")

    today = date.today().isoformat()
    user_focus = (focus_prompt or "").strip()
    focus_block = (
        f"\nCustom focus from user:\n{user_focus}\n"
        if user_focus else
        "\nCustom focus from user:\n(Not provided)\n"
    )

    prompt = f"""
You are an administrative task planner for district governance reviews.

Generate a practical task list from the input text below.
Department: {department_name or "Unknown"}
Source: {source_name}
Today's date: {today}
{focus_block}
Return ONLY valid JSON object, no markdown, no commentary.
JSON schema:
{{
  "suggestions": [
    {{
      "description": "string (required, specific and action-oriented)",
      "assigned_agency": "string or null",
      "priority": "Critical|High|Normal|Low",
      "deadline_date": "YYYY-MM-DD or null",
      "time_given": "string or null",
      "remarks": "short rationale / expected outcome",
      "source_snippet": "short evidence line from input",
      "confidence": 0.0
    }}
  ]
}}

Rules:
- Generate 6 to 20 suggestions.
- Prioritize immediate execution tasks, especially 15-day and 30-day actions where relevant.
- Avoid duplicates and vague items.
- If date is uncertain, use null deadline_date and suggest time_given (like "15 days").
- confidence must be between 0 and 1.

Input text:
{source_text}
""".strip()

    preferred = _normalize_model_name(GEMINI_MODEL)
    model_candidates = [
        preferred,
        "gemini-2.5-flash",
        "gemini-2.5-pro",
        "gemini-2.5-flash-lite",
        "gemini-2.0-flash",
        "gemini-1.5-flash",
    ]
    model_candidates = [_normalize_model_name(m) for m in model_candidates]
    model_candidates = [m for i, m in enumerate(model_candidates) if m and m not in model_candidates[:i]]

    tried = []
    last_error = None
    for model_name in model_candidates:
        tried.append(model_name)
        for output_tokens in _output_token_candidates():
            try:
                raw, _ = _gemini_generate_once(api_key, model_name, prompt, output_tokens)
                parsed = _extract_json_from_text(raw)
                suggestions = parsed.get("suggestions") if isinstance(parsed, dict) else None
                if not isinstance(suggestions, list):
                    raise ValueError("Model response missing 'suggestions' array.")

                cleaned = []
                for i, item in enumerate(suggestions):
                    if not isinstance(item, dict):
                        continue
                    description = str(item.get("description") or "").strip()
                    if len(description) < 6:
                        continue
                    deadline = item.get("deadline_date")
                    if deadline is not None:
                        deadline = str(deadline).strip() or None
                        if deadline:
                            try:
                                date.fromisoformat(deadline)
                            except Exception:
                                deadline = None
                    confidence = item.get("confidence", 0.65)
                    try:
                        confidence = float(confidence)
                    except Exception:
                        confidence = 0.65
                    confidence = max(0.0, min(1.0, confidence))
                    cleaned.append({
                        "id": f"s{i + 1}",
                        "description": description,
                        "assigned_agency": (str(item.get("assigned_agency") or "").strip() or None),
                        "priority": _normalize_priority(item.get("priority")),
                        "deadline_date": deadline,
                        "time_given": (str(item.get("time_given") or "").strip() or None),
                        "remarks": (str(item.get("remarks") or "").strip() or None),
                        "source_snippet": (str(item.get("source_snippet") or "").strip() or None),
                        "confidence": round(confidence, 2),
                    })
                if cleaned:
                    return cleaned[:25]
                raise ValueError("No usable suggestions found in model response.")
            except urllib.error.HTTPError as exc:
                detail = exc.read().decode("utf-8", errors="ignore")
                detail_l = detail.lower()
                if exc.code == 404 and ("not found" in detail_l or "not supported" in detail_l):
                    last_error = GeminiModelNotFoundError(f"Model '{model_name}' unavailable")
                    break
                if exc.code == 400 and (
                    "maxoutputtokens" in detail_l
                    or "max output token" in detail_l
                    or "invalid argument" in detail_l
                    or "unexpected model name format" in detail_l
                ):
                    last_error = RuntimeError(
                        f"Model '{model_name}' rejected request (tokens/model format). Retrying fallback."
                    )
                    continue
                raise RuntimeError(f"Gemini API error ({exc.code}): {detail[:500]}") from exc
            except urllib.error.URLError as exc:
                raise RuntimeError(f"Gemini API network error: {exc}") from exc
            except Exception as exc:
                last_error = exc
                continue

    raise RuntimeError(
        f"No compatible Gemini model found for task suggestions. Tried: {', '.join(tried)}. "
        f"Last error: {last_error}"
    )


def compare_analysis_outputs_with_gemini(
    left_label: str,
    left_analysis: str,
    right_label: str,
    right_analysis: str,
) -> str:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is not set. Please add it in environment variables.")

    left_text = (left_analysis or "").strip()[:16000]
    right_text = (right_analysis or "").strip()[:16000]
    if not left_text or not right_text:
        raise ValueError("Both analyses must contain text to compare.")

    prompt = f"""
You are a governance review comparison analyst.

Compare these two analysis outputs and evaluate progress between dates.

Analysis A (older):
{left_label}

Analysis B (newer):
{right_label}

Return Markdown with exactly these sections:

1) Executive Progress Verdict
- 4-6 bullets summarizing overall movement.

2) Comparable Parameter Matrix
- Markdown table with columns:
  Parameter | Analysis A | Analysis B | Trend (Improved/No Change/Regressed) | Evidence

3) Notable Improvements
- 3-8 specific bullets.

4) Persistent / New Gaps
- 3-8 specific bullets with severity tags (High/Medium/Low).

5) Priority Follow-up Actions
- Markdown table with columns:
  Action | Owner Suggestion | Timeline | Why It Matters

Rules:
- Use only evidence from the provided analyses.
- If evidence is missing, write "Insufficient evidence".
- Keep output practical for department review meetings.

Analysis A text:
{left_text}

Analysis B text:
{right_text}
""".strip()

    preferred = _normalize_model_name(GEMINI_MODEL)
    model_candidates = [
        preferred,
        "gemini-2.5-flash",
        "gemini-2.5-pro",
        "gemini-2.5-flash-lite",
        "gemini-2.0-flash",
        "gemini-1.5-flash",
    ]
    model_candidates = [_normalize_model_name(m) for m in model_candidates]
    model_candidates = [m for i, m in enumerate(model_candidates) if m and m not in model_candidates[:i]]

    tried = []
    last_error = None
    token_candidates = [8192, 6144, 4096, 3072, 2048]

    for model_name in model_candidates:
        tried.append(model_name)
        for output_tokens in token_candidates:
            try:
                result, _ = _gemini_generate_once(api_key, model_name, prompt, output_tokens)
                return result.strip()
            except urllib.error.HTTPError as exc:
                detail = exc.read().decode("utf-8", errors="ignore")
                detail_l = detail.lower()
                if exc.code == 404 and ("not found" in detail_l or "not supported" in detail_l):
                    last_error = GeminiModelNotFoundError(f"Model '{model_name}' unavailable")
                    break
                if exc.code == 400 and (
                    "maxoutputtokens" in detail_l
                    or "max output token" in detail_l
                    or "invalid argument" in detail_l
                    or "unexpected model name format" in detail_l
                ):
                    last_error = RuntimeError(
                        f"Model '{model_name}' rejected comparison request. Retrying fallback."
                    )
                    continue
                raise RuntimeError(f"Gemini API error ({exc.code}): {detail[:500]}") from exc
            except urllib.error.URLError as exc:
                raise RuntimeError(f"Gemini API network error: {exc}") from exc
            except Exception as exc:
                last_error = exc
                continue

    raise RuntimeError(
        f"No compatible Gemini model found for analysis comparison. Tried: {', '.join(tried)}. "
        f"Last error: {last_error}"
    )


def analyze_with_gemini(document_name: str, extracted_text: str, mode: str = "default", custom_prompt: Optional[str] = None) -> str:
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY is not set. Please add it in environment variables.")

    if mode not in {"default", "custom"}:
        raise ValueError("Invalid mode. Use 'default' or 'custom'.")
    if mode == "custom" and not (custom_prompt or "").strip():
        raise ValueError("Custom mode requires a prompt.")

    prompt = _build_default_prompt(document_name, extracted_text)
    if mode == "custom":
        prompt = _build_custom_prompt(document_name, extracted_text, custom_prompt.strip())

    preferred = _normalize_model_name(GEMINI_MODEL)
    model_candidates = [
        preferred,
        "gemini-2.5-flash",
        "gemini-2.5-pro",
        "gemini-2.5-flash-lite",
        "gemini-2.0-flash",
        "gemini-1.5-flash",
    ]
    model_candidates = [_normalize_model_name(m) for m in model_candidates]
    model_candidates = [m for i, m in enumerate(model_candidates) if m and m not in model_candidates[:i]]

    tried = []
    last_error = None

    for model_name in model_candidates:
        tried.append(model_name)
        for output_tokens in _output_token_candidates():
            combined_output = ""
            current_prompt = prompt
            continuation_count = 0
            try:
                while True:
                    chunk, finish_reason = _gemini_generate_once(
                        api_key, model_name, current_prompt, output_tokens
                    )
                    combined_output = f"{combined_output}\n{chunk}".strip() if combined_output else chunk

                    if finish_reason != "MAX_TOKENS":
                        return combined_output.strip()

                    continuation_count += 1
                    if continuation_count > MAX_CONTINUATIONS:
                        return (
                            f"{combined_output.strip()}\n\n"
                            f"_Note: Output hit model token limit after {continuation_count} segments. "
                            f"Re-run with Custom mode for a narrower prompt if needed._"
                        )

                    current_prompt = _build_continuation_prompt(combined_output)
            except urllib.error.HTTPError as exc:
                detail = exc.read().decode("utf-8", errors="ignore")
                detail_l = detail.lower()
                if exc.code == 404 and ("not found" in detail_l or "not supported" in detail_l):
                    last_error = GeminiModelNotFoundError(f"Model '{model_name}' unavailable")
                    break
                if exc.code == 400 and (
                    "maxoutputtokens" in detail_l
                    or "max output token" in detail_l
                    or "invalid argument" in detail_l
                    or "unexpected model name format" in detail_l
                ):
                    last_error = RuntimeError(
                        f"Model '{model_name}' rejected maxOutputTokens={output_tokens}. Retrying with lower limit."
                    )
                    continue
                raise RuntimeError(f"Gemini API error ({exc.code}): {detail[:500]}") from exc
            except urllib.error.URLError as exc:
                raise RuntimeError(f"Gemini API network error: {exc}") from exc
            except Exception as exc:
                last_error = exc
                continue

    raise RuntimeError(
        f"No compatible Gemini model found. Tried: {', '.join(tried)}. "
        f"Set GEMINI_MODEL in env to a model from ListModels. Last error: {last_error}"
    )
